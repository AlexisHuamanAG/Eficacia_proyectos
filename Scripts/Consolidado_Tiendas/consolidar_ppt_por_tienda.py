from __future__ import annotations

import argparse
import logging
import re
from collections import OrderedDict, defaultdict
from pathlib import Path

from pptx import Presentation


DEFAULT_FOLDER = Path(
    r"C:\1\OneDrive - Eficacia\KENVUE - Front\EXHIBICIONES ACUMULADAS Y ELEMENTOS\2026\5. Mayo\HS NEGOCIADA"
)
DEFAULT_MODEL = Path(
    r"C:\1\OneDrive - Eficacia\KENVUE - Front\EXHIBICIONES ACUMULADAS Y ELEMENTOS\2026\Scripts\MODELO.pptx"
)

STORE_ORDER = ["Metro", "Tottus", "Plaza Vea", "Wong", "Vivanda", "Inka Farma", "Mifarma", "Aruma"]

CODE_STORE_PREFIXES = {
    "TP": "Tottus",
    "T": "Wong",
    "H": "Metro",
    "S": "Metro",
    "PV": "Plaza Vea",
    "P": "Plaza Vea",
    "IF": "Inka Farma",
    "INK": "Inka Farma",
    "MF": "Mifarma",
    "MIF": "Mifarma",
    "V": "Vivanda",
    "VIV": "Vivanda",
    "AR": "Aruma",
}


def clean_title(text: str) -> str:
    return " ".join(text.replace("\n", " ").split()).strip()


def clean_element_name(stem: str) -> str:
    name = stem.strip()

    patterns = [
        r"\s*[-_]\s*\d{1,3}$",
        r"\s*\(\s*\d{1,3}\s*\)$",
        r"\s+(?:PARTE|PART|PT)\s*\d{1,3}$",
        r"\s*[-_]\s*(?:PARTE|PART|PT)\s*\d{1,3}$",
        r"\s+(?:PTE)\s*\d{1,3}$",
    ]

    changed = True
    while changed:
        changed = False
        for pattern in patterns:
            cleaned = re.sub(pattern, "", name, flags=re.IGNORECASE).strip()
            if cleaned != name:
                name = cleaned
                changed = True

    return clean_title(name).title()


def title_case_name(text: str) -> str:
    return clean_element_name(text)


def iter_pptx_files(folder: Path):
    for path in sorted(folder.rglob("*.pptx")):
        if path.name.startswith("~$"):
            continue
        if path.name.upper().startswith("MODELO"):
            continue
        if path.name.upper().startswith("CONSOLIDADO"):
            continue
        yield path


def all_slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = clean_title(shape.text_frame.text)
            if text:
                parts.append(text)
    return " | ".join(parts)


def top_text_candidates(slide):
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = clean_title(shape.text_frame.text)
        if not text:
            continue
        candidates.append((int(shape.top), int(shape.left), text))
    return sorted(candidates)


def extract_store_from_pdv(pdv: str) -> str | None:
    normalized = clean_title(pdv).upper()
    if normalized.startswith("PLAZA VEA"):
        return "Plaza Vea"
    if normalized.startswith("TOTTUS"):
        return "Tottus"
    if normalized.startswith("METRO"):
        return "Metro"
    if normalized.startswith("WONG"):
        return "Wong"
    if normalized.startswith("VIVANDA"):
        return "Vivanda"
    if normalized.startswith("INKA FARMA") or normalized.startswith("INKAFARMA"):
        return "Inka Farma"
    if normalized.startswith("MIFARMA") or normalized.startswith("MI FARMA"):
        return "Mifarma"
    if normalized.startswith("ARUMA"):
        return "Aruma"
    return None


def extract_store_from_code(text: str) -> str | None:
    match = re.search(r"C[oó]digo\s+del\s+PDV:\s*([A-Z]+)\d+", text, flags=re.IGNORECASE)
    if not match:
        return None

    code_prefix = match.group(1).upper()
    for prefix in sorted(CODE_STORE_PREFIXES, key=len, reverse=True):
        if code_prefix.startswith(prefix):
            return CODE_STORE_PREFIXES[prefix]
    return None


def extract_pdv_and_store(slide, element: str) -> tuple[str | None, str | None]:
    element_norm = clean_title(element).upper()
    for _, _, text in top_text_candidates(slide):
        text_norm = text.upper()
        if text_norm == element_norm:
            continue
        if text_norm.startswith("INFORME DE FOTOS"):
            continue
        if text_norm.startswith("FOTO DE"):
            continue

        store = extract_store_from_pdv(text)
        if store:
            return text, store

    text = all_slide_text(slide)
    return None, extract_store_from_code(text)


def build_manifest(folder: Path):
    manifest = []
    skipped = []

    for path in iter_pptx_files(folder):
        element = title_case_name(path.stem)
        prs = Presentation(str(path))

        for slide_no, slide in enumerate(prs.slides, 1):
            pdv, store = extract_pdv_and_store(slide, element)
            if not store:
                skipped.append((path, slide_no, "No se pudo identificar tienda"))
                continue

            manifest.append(
                {
                    "source": path,
                    "slide_no": slide_no,
                    "store": store,
                    "element": element,
                    "pdv": pdv or "",
                }
            )

    return manifest, skipped


def grouped_manifest(manifest):
    grouped = defaultdict(lambda: defaultdict(list))
    for item in manifest:
        grouped[item["store"]][item["element"]].append(item)
    return grouped


def store_sort_key(store: str):
    if store in STORE_ORDER:
        return (0, STORE_ORDER.index(store))
    return (1, store)


def find_template_path(folder: Path, model_path: Path | None = None) -> Path | None:
    if model_path and model_path.exists():
        return model_path
    if DEFAULT_MODEL.exists():
        return DEFAULT_MODEL
    candidates = sorted(folder.glob("*MODELO*.pptx"))
    return candidates[0] if candidates else None


def replace_first_text(slide, text: str):
    for shape in slide.Shapes:
        if not shape.HasTextFrame:
            continue
        if not shape.TextFrame.HasText:
            continue
        shape.TextFrame.TextRange.Text = text
        shape.TextFrame.TextRange.Font.Size = 40
        return True
    return False


def add_divider_slide_from_model(presentation, model_presentation, title: str):
    model_presentation.Slides(1).Copy()
    pasted = presentation.Slides.Paste(presentation.Slides.Count + 1)
    slide = pasted.Item(1)
    replace_first_text(slide, title)
    return slide


def ensure_element_header(slide, element: str):
    page_width = slide.Parent.PageSetup.SlideWidth
    page_height = slide.Parent.PageSetup.SlideHeight

    for shape in slide.Shapes:
        if not shape.HasTextFrame:
            continue
        if not shape.TextFrame.HasText:
            continue

        text = shape.TextFrame.TextRange.Text.strip()
        if text == element and shape.Left <= page_width * 0.18 and shape.Top <= page_height * 0.14:
            shape.TextFrame.TextRange.Font.Name = "Arial"
            shape.TextFrame.TextRange.Font.Size = 12
            return shape

    mso_text_orientation_horizontal = 1
    shape = slide.Shapes.AddTextbox(
        mso_text_orientation_horizontal,
        28,
        18,
        180,
        24,
    )
    text_range = shape.TextFrame.TextRange
    text_range.Text = element
    text_range.Font.Name = "Arial"
    text_range.Font.Size = 12
    text_range.Font.Bold = True
    return shape


def create_consolidated_pptx(folder: Path, output_path: Path, manifest, model_path_arg: Path | None = None) -> None:
    import win32com.client

    grouped = grouped_manifest(manifest)
    model_path = find_template_path(folder, model_path_arg)
    if model_path is None:
        raise FileNotFoundError(f"No se encontro MODELO.pptx en {DEFAULT_MODEL.parent} ni en {folder}")

    app = win32com.client.DispatchEx("PowerPoint.Application")
    try:
        app.DisplayAlerts = 0
    except Exception:
        pass

    target = None
    model_presentation = None

    try:
        model_presentation = app.Presentations.Open(str(model_path), True, False, False)
        target = app.Presentations.Add(0)
        target.ApplyTemplate(str(model_path))
        target.Slides.Add(1, 12)

        for store in sorted(grouped, key=store_sort_key):
            add_divider_slide_from_model(target, model_presentation, store)

            for element in sorted(grouped[store]):
                items = grouped[store][element]
                add_divider_slide_from_model(target, model_presentation, element)

                for item in items:
                    target.Slides.InsertFromFile(
                        str(item["source"]),
                        target.Slides.Count,
                        item["slide_no"],
                        item["slide_no"],
                    )
                    ensure_element_header(target.Slides(target.Slides.Count), item["element"])

        target.Slides(1).Delete()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        if output_path.exists():
            output_path.unlink()
        target.SaveAs(str(output_path))
    finally:
        if target is not None:
            target.Close()
        if model_presentation is not None:
            model_presentation.Close()
        try:
            app.Quit()
        except Exception:
            pass


def write_log(output_log: Path, manifest, skipped) -> None:
    grouped = grouped_manifest(manifest)
    lines = []
    lines.append(f"Diapositivas consolidadas: {len(manifest)}")
    lines.append(f"Diapositivas omitidas: {len(skipped)}")
    lines.append("")

    for store in sorted(grouped, key=store_sort_key):
        store_total = sum(len(items) for items in grouped[store].values())
        lines.append(f"{store}: {store_total}")
        for element in sorted(grouped[store]):
            lines.append(f"  - {element}: {len(grouped[store][element])}")
        lines.append("")

    if skipped:
        lines.append("Omitidas:")
        for path, slide_no, reason in skipped:
            lines.append(f"  - {path.name} diapositiva {slide_no}: {reason}")

    output_log.write_text("\n".join(lines), encoding="utf-8")


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Consolida PPTX agrupando diapositivas por tienda y elemento."
    )
    parser.add_argument("folder", nargs="?", default=str(DEFAULT_FOLDER), help="Carpeta raiz con PPTX.")
    parser.add_argument("--output", help="Ruta del PPTX consolidado.")
    parser.add_argument(
        "--model",
        default=str(DEFAULT_MODEL),
        help="Ruta de MODELO.pptx para separadores y tema visual.",
    )
    parser.add_argument("--dry-run", action="store_true", help="Solo analiza y muestra resumen.")
    args = parser.parse_args()

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    folder = Path(args.folder)
    if not folder.exists() or not folder.is_dir():
        logging.error("La carpeta no existe o no es una carpeta: %s", folder)
        return 1

    output_path = Path(args.output) if args.output else folder / f"CONSOLIDADO_{folder.name}.pptx"
    output_log = output_path.with_suffix(".log.txt")

    manifest, skipped = build_manifest(folder)
    if not manifest:
        logging.error("No se encontraron diapositivas con tienda identificable.")
        write_log(output_log, manifest, skipped)
        return 1

    write_log(output_log, manifest, skipped)
    logging.info("Diapositivas a consolidar: %s", len(manifest))
    logging.info("Diapositivas omitidas: %s", len(skipped))
    logging.info("Log: %s", output_log)

    grouped = grouped_manifest(manifest)
    for store in sorted(grouped, key=store_sort_key):
        total = sum(len(items) for items in grouped[store].values())
        logging.info("%s: %s diapositivas", store, total)

    if args.dry_run:
        logging.info("Modo dry-run: no se genero PPTX.")
        return 0

    create_consolidated_pptx(folder, output_path, manifest, Path(args.model))
    logging.info("PPTX consolidado generado: %s", output_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
