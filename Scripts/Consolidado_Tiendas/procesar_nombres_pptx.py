from __future__ import annotations

import argparse
import logging
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_FOLDER = Path(
    r"C:\1\OneDrive - Eficacia\KENVUE - Front\EXHIBICIONES ACUMULADAS Y ELEMENTOS\2026\5. Mayo\HS NEGOCIADA"
)


def clean_element_name(stem: str) -> str:
    name = stem.strip()

    patterns = [
        r"\s*[-_]\s*\d{1,3}$",
        r"\s*\(\s*\d{1,3}\s*\)$",
        r"\s+(?:PARTE|PART|PT)\s*\d{1,3}$",
        r"\s*[-_]\s*(?:PARTE|PART|PT)\s*\d{1,3}$",
        r"\s+(?:PTE)\s*\d{1,3}$",
    ]

    changed = True
    while changed:
        changed = False
        for pattern in patterns:
            cleaned = re.sub(pattern, "", name, flags=re.IGNORECASE).strip()
            if cleaned != name:
                name = cleaned
                changed = True

    return " ".join(name.split()).title()


def display_name_from_file(path: Path) -> str:
    """Convierte 'MUEBLE LATERAL - 002.pptx' en 'Mueble Lateral'."""
    return clean_element_name(path.stem)


def iter_candidate_shapes(slide, slide_width: int, slide_height: int):
    max_left = slide_width * 0.16
    max_top = slide_height * 0.12

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        text = shape.text_frame.text.strip()
        if not text:
            continue

        if shape.left <= max_left and shape.top <= max_top:
            # Prioriza lo que este mas arriba e izquierda; despues cuadros mas pequenos.
            score = (
                int(shape.top) * 10
                + int(shape.left) * 5
                + int(shape.width) // 1000
                + int(shape.height) // 1000
            )
            yield score, shape


def find_top_left_textbox(slide, slide_width: int, slide_height: int):
    candidates = sorted(iter_candidate_shapes(slide, slide_width, slide_height), key=lambda item: item[0])
    return candidates[0][1] if candidates else None


def find_filename_header(slide, slide_width: int, slide_height: int, replacement: str):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        if shape.left > slide_width * 0.16 or shape.top > slide_height * 0.12:
            continue

        text = shape.text_frame.text.strip()
        if text == replacement:
            return shape

    return None


def set_header_font(shape) -> None:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(12)


def copy_font(source_run, target_run) -> None:
    source_font = source_run.font
    target_font = target_run.font

    target_font.name = source_font.name
    target_font.size = source_font.size
    target_font.bold = source_font.bold
    target_font.italic = source_font.italic
    target_font.underline = source_font.underline

    if source_font.color.type is not None:
        if getattr(source_font.color, "rgb", None) is not None:
            target_font.color.rgb = source_font.color.rgb


def first_run(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            return run
    return None


def add_textbox_like(slide, template_shape, new_text: str):
    shape = slide.shapes.add_textbox(
        template_shape.left,
        template_shape.top,
        template_shape.width,
        template_shape.height,
    )
    shape.rotation = template_shape.rotation

    text_frame = shape.text_frame
    text_frame.margin_left = template_shape.text_frame.margin_left
    text_frame.margin_right = template_shape.text_frame.margin_right
    text_frame.margin_top = template_shape.text_frame.margin_top
    text_frame.margin_bottom = template_shape.text_frame.margin_bottom
    text_frame.word_wrap = template_shape.text_frame.word_wrap
    text_frame.auto_size = template_shape.text_frame.auto_size

    paragraph = text_frame.paragraphs[0]
    template_paragraph = template_shape.text_frame.paragraphs[0]
    paragraph.alignment = template_paragraph.alignment or PP_ALIGN.LEFT
    paragraph.level = template_paragraph.level

    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.text = new_text

    template_run = first_run(template_shape)
    if template_run is not None:
        copy_font(template_run, run)

    return shape


def add_default_filename_header(slide, new_text: str):
    shape = slide.shapes.add_textbox(Inches(0.35), Inches(0.20), Inches(1.6), Inches(0.25))
    text_frame = shape.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    run.text = new_text
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.bold = True
    return shape


def replace_text_preserve_format(shape, new_text: str) -> None:
    text_frame = shape.text_frame

    if not text_frame.paragraphs:
        text_frame.text = new_text
        return

    first_paragraph = text_frame.paragraphs[0]

    if not first_paragraph.runs:
        first_paragraph.text = new_text
    else:
        first_paragraph.runs[0].text = new_text

        for run in first_paragraph.runs[1:]:
            run.text = ""

    for paragraph in text_frame.paragraphs[1:]:
        for run in paragraph.runs:
            run.text = ""


def process_presentation(path: Path, dry_run: bool = False) -> tuple[int, int]:
    prs = Presentation(str(path))
    replacement = display_name_from_file(path)
    changed = 0
    missing = 0

    for slide in prs.slides:
        shape = find_filename_header(slide, prs.slide_width, prs.slide_height, replacement)
        if shape is None:
            changed += 1
            if not dry_run:
                add_default_filename_header(slide, replacement)
        else:
            changed += 1
            if not dry_run:
                set_header_font(shape)

    if changed and not dry_run:
        prs.save(str(path))

    return changed, missing


def iter_pptx_files(folder: Path):
    for path in folder.rglob("*.pptx"):
        if path.name.startswith("~$"):
            continue
        yield path


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Reemplaza el texto superior izquierdo de todas las diapositivas por el nombre del archivo PPTX."
    )
    parser.add_argument(
        "folder",
        nargs="?",
        default=str(DEFAULT_FOLDER),
        help="Carpeta raiz a procesar. Si se omite, usa la ruta configurada en el script.",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Muestra que cambiaria sin guardar los archivos.",
    )
    args = parser.parse_args()

    folder = Path(args.folder)
    if not folder.exists() or not folder.is_dir():
        logging.error("La carpeta no existe o no es una carpeta: %s", folder)
        return 1

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    files = list(iter_pptx_files(folder))
    if not files:
        logging.warning("No se encontraron archivos .pptx en: %s", folder)
        return 0

    logging.info("Archivos .pptx encontrados: %s", len(files))

    total_changed = 0
    total_missing = 0

    for path in files:
        try:
            changed, missing = process_presentation(path, dry_run=args.dry_run)
            total_changed += changed
            total_missing += missing
            logging.info(
                "Procesado: %s | diapositivas actualizadas: %s | sin cuadro detectado: %s",
                path,
                changed,
                missing,
            )
        except Exception as exc:
            logging.exception("Error procesando %s: %s", path, exc)

    if args.dry_run:
        logging.info("Modo dry-run: no se guardaron cambios.")

    logging.info("Total de diapositivas actualizadas: %s", total_changed)
    logging.info("Total de diapositivas sin cuadro detectado: %s", total_missing)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
